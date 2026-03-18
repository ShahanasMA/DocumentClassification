from flask import Flask, render_template, request, redirect, session, send_from_directory
import sqlite3, joblib, os
from werkzeug.security import generate_password_hash, check_password_hash
from werkzeug.utils import secure_filename
from datetime import datetime
from PyPDF2 import PdfReader
from pptx import Presentation
from database import init_db

app = Flask(__name__)
app.secret_key = "secret"

init_db()

model = joblib.load("document_classifier.pkl")

BASE = "documents"
os.makedirs(BASE, exist_ok=True)

DB = "database.db"

ALL_CATEGORIES = ["HR", "Finance", "Legal", "Administration"]


def get_db():
    return sqlite3.connect(DB)


# 🔔 ---------------- NOTIFICATIONS ----------------
@app.context_processor
def inject_notifications():

    notifications = []

    if "user" in session and session.get("role") == "admin":

        conn = get_db()
        cur = conn.cursor()

        cur.execute("""
        SELECT d.filename
        FROM review_requests r
        JOIN documents d ON d.id = r.document_id
        WHERE r.status='pending'
        ORDER BY r.id DESC
        LIMIT 5
        """)
        docs = cur.fetchall()

        for d in docs:
            notifications.append(f"📄 Document '{d[0]}' needs review")

        cur.execute("""
        SELECT username FROM users
        WHERE status='pending' AND role='user'
        ORDER BY id DESC
        LIMIT 5
        """)
        users = cur.fetchall()

        for u in users:
            notifications.append(f"👤 User '{u[0]}' waiting for approval")

        conn.close()

    return dict(notifications=notifications, notif_count=len(notifications))


# ---------------- TEXT EXTRACTION ----------------
def extract(file, name):

    if name.endswith(".txt"):
        return file.read().decode("utf-8", errors="ignore")

    if name.endswith(".pdf"):
        reader = PdfReader(file)
        return "".join(page.extract_text() or "" for page in reader.pages)

    if name.endswith(".pptx"):
        prs = Presentation(file)
        text = ""
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text
        return text

    return ""


# ---------------- REGISTER ----------------
@app.route("/register", methods=["GET", "POST"])
def register():

    if request.method == "POST":

        if request.form["password"] != request.form["confirm"]:
            return render_template("register.html", msg="Passwords do not match")

        conn = get_db()
        cur = conn.cursor()

        try:
            cur.execute("""
            INSERT INTO users(name,email,username,password,role,status,department)
            VALUES(?,?,?,?,?,?,?)
            """, (
                request.form["name"],
                request.form["email"],
                request.form["username"],
                generate_password_hash(request.form["password"]),
                "user",
                "pending",
                request.form["department"]
            ))

            conn.commit()
            msg = "Registration successful. Wait for admin approval."

        except sqlite3.IntegrityError:
            msg = "Username already exists"

        conn.close()
        return render_template("register.html", msg=msg)

    return render_template("register.html")


# ---------------- LOGIN ----------------
@app.route("/login", methods=["GET", "POST"])
def login():

    if request.method == "POST":

        conn = get_db()
        cur = conn.cursor()

        cur.execute(
            "SELECT password,role,status,department FROM users WHERE username=?",
            (request.form["username"],)
        )

        user = cur.fetchone()
        conn.close()

        if user and check_password_hash(user[0], request.form["password"]):

            if user[2] != "approved":
                return render_template("login.html", error="Account waiting for admin approval")

            session["user"] = request.form["username"]
            session["role"] = user[1]
            session["department"] = user[3]
            return redirect("/")

        return render_template("login.html", error="Invalid login")

    return render_template("login.html")


# ---------------- PROFILE ----------------
@app.route("/profile")
def profile():

    if "user" not in session:
        return redirect("/login")

    conn = get_db()
    cur = conn.cursor()

    cur.execute("""
    SELECT name,email,username,department,role
    FROM users WHERE username=?
    """, (session["user"],))

    user = cur.fetchone()
    conn.close()

    return render_template("profile.html", user=user)


# ---------------- HOME ----------------
@app.route("/")
def home():
    if "user" not in session:
        return redirect("/login")
    return render_template("upload.html")


# ---------------- UPLOAD ----------------
@app.route("/", methods=["POST"])
def upload():

    file = request.files["file"]
    filename = secure_filename(file.filename)

    if not filename.endswith((".txt", ".pdf", ".pptx")):
        return "Invalid file type"

    text = extract(file, filename.lower())

    category = model.predict([text])[0]
    confidence = float(max(model.predict_proba([text])[0]))

    percent = int(confidence * 100)
    color = "green" if confidence >= 0.7 else "yellow" if confidence >= 0.4 else "red"

    path = os.path.join(BASE, category)
    os.makedirs(path, exist_ok=True)

    file.seek(0)
    with open(os.path.join(path, filename), "wb") as f:
        f.write(file.read())

    conn = get_db()
    cur = conn.cursor()

    cur.execute("""
    INSERT INTO documents(filename,category,confidence,preview,uploaded_by,upload_time,status)
    VALUES(?,?,?,?,?,?,?)
    """, (
        filename,
        category,
        confidence,
        text[:150],
        session["user"],
        datetime.now(),
        "classified"
    ))

    doc_id = cur.lastrowid

    if confidence < 0.40:
        cur.execute("""
        INSERT INTO review_requests(document_id,message,status)
        VALUES(?,?,?)
        """, (doc_id, "Low confidence auto-flagged", "pending"))

    conn.commit()
    conn.close()

    return render_template("upload.html",
                           category=category,
                           percent=percent,
                           color=color,
                           filename=filename)


# ---------------- REQUEST REVIEW ----------------
@app.route("/request_review", methods=["POST"])
def request_review():

    filename = request.form["filename"]
    message = request.form["message"]

    conn = get_db()
    cur = conn.cursor()

    cur.execute("""
    SELECT id FROM documents
    WHERE filename=? ORDER BY id DESC LIMIT 1
    """, (filename,))

    doc = cur.fetchone()

    if doc:
        cur.execute("""
        INSERT INTO review_requests(document_id,message,status)
        VALUES(?,?,?)
        """, (doc[0], message, "pending"))

        conn.commit()

    conn.close()
    return redirect("/")


# ---------------- VIEW FILE ----------------
@app.route("/view/<category>/<filename>")
def view_file(category, filename):

    if "user" not in session:
        return redirect("/login")

    # ✅ NEW LINE (DO NOT REMOVE ANYTHING ELSE)
    source = request.args.get("source")

    return render_template("view.html",
                           category=category,
                           filename=filename,
                           source=source)

# ---------------- CATEGORIES ----------------
@app.route("/categories")
def categories():

    if session["role"] == "admin":
        cats = ALL_CATEGORIES
    else:
        cats = [session["department"]]

    return render_template("categories.html", cats=cats)


# ---------------- CATEGORY VIEW ----------------
@app.route("/category/<cat>")
def category(cat):

    conn = get_db()
    cur = conn.cursor()

    cur.execute("""
    SELECT filename,confidence,preview,upload_time
    FROM documents WHERE category=?
    """, (cat,))

    rows = cur.fetchall()

    files = []
    for r in rows:
        percent = int(r[1] * 100)
        color = "green" if r[1] >= 0.7 else "yellow" if r[1] >= 0.4 else "red"

        files.append({
            "filename": r[0],
            "percent": percent,
            "color": color,
            "preview": r[2],
            "time": r[3]
        })

    conn.close()

    return render_template("category_view.html", files=files, category=cat)


# ---------------- ADMIN PANEL ----------------
@app.route("/admin")
def admin():

    conn = get_db()
    cur = conn.cursor()

    cur.execute("""
    SELECT d.id,d.filename,d.category,r.message,r.status
    FROM review_requests r
    JOIN documents d ON d.id=r.document_id
    ORDER BY CASE WHEN r.status='pending' THEN 0 ELSE 1 END, r.id DESC
    """)

    rows = cur.fetchall()

    data = [{
        "id": r[0],
        "filename": r[1],
        "category": r[2],
        "message": r[3],
        "status": r[4]
    } for r in rows]

    conn.close()

    return render_template("admin.html", data=data)


# ---------------- ADMIN UPDATE ----------------
@app.route("/admin/update", methods=["POST"])
def update_category():

    conn = get_db()
    cur = conn.cursor()

    cur.execute("UPDATE documents SET category=? WHERE id=?",
                (request.form["new_category"], request.form["doc_id"]))

    cur.execute("UPDATE review_requests SET status='resolved' WHERE document_id=?",
                (request.form["doc_id"],))

    conn.commit()
    conn.close()

    return redirect("/admin")

# ---------------- USER APPROVAL ----------------
@app.route("/admin/users")
def admin_users():

    if session.get("role") != "admin":
        return "Unauthorized"

    conn = get_db()
    cur = conn.cursor()

    cur.execute("""
    SELECT id,name,username,status,department,email
    FROM users
    WHERE role='user'
    ORDER BY 
        CASE WHEN status='pending' THEN 0 ELSE 1 END,
        id DESC
    """)

    users = cur.fetchall()
    conn.close()

    return render_template("admin_users.html", users=users)


# ---------------- APPROVE USER ----------------
@app.route("/admin/approve/<int:id>")
def approve_user(id):

    if session.get("role") != "admin":
        return "Unauthorized"

    conn = get_db()
    cur = conn.cursor()

    cur.execute("UPDATE users SET status='approved' WHERE id=?", (id,))

    conn.commit()
    conn.close()

    return redirect("/admin/users")


# ---------------- REJECT USER ----------------
@app.route("/admin/reject/<int:id>")
def reject_user(id):

    if session.get("role") != "admin":
        return "Unauthorized"

    conn = get_db()
    cur = conn.cursor()

    cur.execute("DELETE FROM users WHERE id=?", (id,))

    conn.commit()
    conn.close()

    return redirect("/admin/users")


# ---------------- FILE ACCESS ----------------
@app.route("/open/<category>/<path:filename>")
def open_file(category, filename):

    return send_from_directory(os.path.join(BASE, category), filename)


@app.route("/download/<category>/<path:filename>")
def download_file(category, filename):

    return send_from_directory(os.path.join(BASE, category), filename, as_attachment=True)


# ---------------- LOGOUT ----------------
@app.route("/logout")
def logout():
    session.clear()
    return redirect("/login")


if __name__ == "__main__":
    app.run(debug=True)
